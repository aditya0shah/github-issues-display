from github import Github
import pptx
from pptx.util import Pt, Inches


# auth = Auth.Token("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")
g = Github()
prs = pptx.Presentation("public/slide.pptx")
repo = g.get_repo("FRC-1294/frc2024")

slide = prs.slides[0]




issues = repo.get_issues(labels=["Mechanics","High Priority"])

print(issues)

for issue in issues:
    print (issue.title)


# for shape in slide.shapes:
#     if shape.has_text_frame:
#         text_frame = shape.text_frame
#         for paragraph in text_frame.paragraphs:
#             for run in paragraph.runs:
#                 text = run.text
#                 if "Task" in text:
#                     new_text = text.replace("Task", "new Task")
#                     run.text = new_text

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(Inches(3.43), Inches(1.86), Inches(2.02), Inches(0.91))
tf = txBox.text_frame
temp = tf.add_paragraph()
temp.text = issues[2].title
temp.font.size = Pt(20)

try:
    p = tf.add_paragraph()
    p.text = issues[2].assignee
    p.font.size = Pt(30)
    prs.save("public/newSlide.pptx")
except:
    prs.save("public/newSlide.pptx")



# slide.placeholders[1].text = issues[0].title
# print("thing, ", slide.placeholders[0].text)
#prs.save("public/newSlide.pptx")


    